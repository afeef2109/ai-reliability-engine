from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.colors import Color, HexColor
from reportlab.lib.pagesizes import landscape
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

ROOT = Path(r"C:\Users\AFEEF M\OneDrive\Documents\Playground")
OUT_DIR = ROOT / "deck"
ASSET_DIR = OUT_DIR / "assets"
ASSET_DIR.mkdir(parents=True, exist_ok=True)
LOGO = ROOT / "public" / "brand" / "stable-logo.png"
PRODUCT_IMG = ASSET_DIR / "product-snapshot.png"
PDF_PATH = OUT_DIR / "stable-investor-deck.pdf"
PPTX_PATH = OUT_DIR / "stable-investor-deck.pptx"

BG = "#060606"
CARD = "#111111"
CARD_ALT = "#191312"
ACCENT = "#FF5C3A"
ACCENT_SOFT = "#FF9A73"
TEXT = "#FFF4EF"
MUTED = "#D1B2A4"
BORDER = "#3A2A27"
GRID_RGBA = (255, 255, 255, 12)

slides = [
    {
        "tag": "AI RELIABILITY ENGINE",
        "title": "STABLE",
        "subtitle": "Detect wrong, risky, and unreliable AI outputs before they reach users.",
        "bullets": [
            "Reliability infrastructure for production AI systems",
            "Highlights risky text, explains failures, and scores output confidence"
        ],
    },
    {
        "tag": "THE PROBLEM",
        "title": "AI is shipping faster than teams can trust it.",
        "subtitle": "LLMs are moving into customer workflows, but failure modes remain invisible until they cause damage.",
        "bullets": [
            "Hallucinated facts sound plausible",
            "Unsafe advice can reach end users",
            "Wrong answers erode trust, increase legal risk, and slow adoption"
        ],
    },
    {
        "tag": "WHY NOW",
        "title": "The market has moved from model access to model trust.",
        "subtitle": "Every company is becoming an AI product company. Reliability is now the bottleneck to deployment.",
        "bullets": [
            "AI is entering support, operations, healthcare, legal, and finance",
            "Enterprises now demand auditability and guardrails",
            "Teams need infrastructure, not just prompts"
        ],
    },
    {
        "tag": "OUR SOLUTION",
        "title": "STABLE is a reliability layer between model output and real-world use.",
        "subtitle": "We generate or ingest an AI response, then run a second-pass evaluator to detect risky or unreliable content.",
        "bullets": [
            "Flags hallucinations, unsafe output, weak factual grounding, and misleading certainty",
            "Highlights the exact text that may be wrong",
            "Returns confidence score, risk level, and explanation in structured form"
        ],
    },
    {
        "tag": "PRODUCT",
        "title": "The product makes AI failure visible in seconds.",
        "subtitle": "Instead of pass or fail, STABLE shows exactly what broke and why.",
        "bullets": [
            "Prompt in, response out, evaluator on top",
            "Highlighted risky phrase inside the model answer",
            "Clear confidence score, risk badge, issues, and explanation"
        ],
        "image": str(PRODUCT_IMG),
    },
    {
        "tag": "WHAT WE DETECT",
        "title": "Failure modes we detect",
        "subtitle": "Built for the real behaviors that break trust in production AI.",
        "bullets": [
            "Hallucinations",
            "Incorrect facts",
            "Unsafe outputs",
            "Inconsistent answers",
            "Low-confidence reasoning",
            "Weak or unverifiable claims"
        ],
    },
    {
        "tag": "TRACTION",
        "title": "Early traction story: strong signal, fast demos, clear buyer pain.",
        "subtitle": "This is still early, but the product already demonstrates a visible and urgent category need.",
        "bullets": [
            "Live interactive demo built and deployable on Vercel",
            "Core reliability pipeline running through OpenRouter with second-pass analysis",
            "Design and workflow already strong enough for pilots, investor meetings, and customer discovery"
        ],
        "stats": [
            ("Demo-ready", "Investor and customer presentation quality"),
            ("2-pass engine", "Generation plus evaluator architecture live"),
            ("Clear wedge", "AI startups and enterprise copilots")
        ]
    },
    {
        "tag": "MARKET",
        "title": "A growing reliability layer inside the AI stack.",
        "subtitle": "As LLM usage expands, every serious deployment needs trust, visibility, and control.",
        "bullets": [
            "TAM: global AI software and infrastructure budgets expanding rapidly",
            "SAM: teams deploying copilots, agents, and workflow automation with external-facing risk",
            "SOM: early wedge in startups and mid-market teams adopting AI into operations"
        ],
        "stats": [
            ("Category", "AI evals, observability, guardrails"),
            ("Buyer", "CTO, Head of AI, Product, Risk, Compliance"),
            ("Why now", "Trust is the bottleneck to production AI")
        ]
    },
    {
        "tag": "CUSTOMERS",
        "title": "Initial wedge: teams already shipping AI into real workflows.",
        "subtitle": "If an AI response reaches a user, that product needs a trust layer.",
        "bullets": [
            "AI startups and agent builders",
            "Enterprise copilots and support automation platforms",
            "Healthcare, legal, fintech, and compliance-heavy AI products"
        ],
    },
    {
        "tag": "BUSINESS MODEL",
        "title": "Infrastructure economics with recurring expansion paths.",
        "subtitle": "STABLE can monetize as a SaaS plus usage-based reliability platform.",
        "bullets": [
            "Subscription plus metered pricing per response analyzed",
            "Enterprise tiers for dashboards, alerts, policy customization, and audit logs",
            "Expansion through monitoring, eval pipelines, and deployment gating"
        ],
    },
    {
        "tag": "TEAM",
        "title": "Lean team, sharp problem focus, product-first execution.",
        "subtitle": "The founding advantage is speed: turning a painful and visible AI problem into usable infrastructure fast.",
        "bullets": [
            "Founder-led product and engineering execution",
            "Focus on design-forward infrastructure that demos well and solves real risk",
            "Next hires: full-stack engineering, AI evals, and enterprise GTM"
        ],
        "stats": [
            ("Founder DNA", "Product + engineering + execution velocity"),
            ("Next hires", "Reliability engineer, frontend/product design, GTM"),
            ("Operating style", "Fast iteration with enterprise positioning")
        ]
    },
    {
        "tag": "WHY WE WIN",
        "title": "We are not another model. We are the trust layer.",
        "subtitle": "STABLE sits in the reliability layer of the AI stack, where urgency and value both compound.",
        "bullets": [
            "Clear pain with visible ROI: reduce brand risk and ship faster",
            "Product value is instantly understandable in demos",
            "Positioned as AI observability plus QA plus reliability infrastructure"
        ],
    },
    {
        "tag": "THE ASK",
        "title": "We are building the system companies use to trust AI in production.",
        "subtitle": "Funding accelerates product depth, enterprise readiness, and pilot customer growth.",
        "bullets": [
            "Use of funds: product, engineering, reliability engine, and GTM",
            "Goal: become the default checkpoint before AI reaches users",
            "Closing line: Before AI reaches users, it should pass through STABLE"
        ],
    },
]


def load_font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\Arialbd.ttf" if bold else r"C:\Windows\Fonts\Arial.ttf",
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


def rounded(draw, box, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def create_product_snapshot():
    w, h = 1600, 980
    image = Image.new("RGBA", (w, h), BG)
    draw = ImageDraw.Draw(image, "RGBA")

    for x in range(0, w, 60):
        draw.line((x, 0, x, h), fill=GRID_RGBA, width=1)
    for y in range(0, h, 60):
        draw.line((0, y, w, y), fill=GRID_RGBA, width=1)

    draw.ellipse((-80, -40, 420, 420), fill=(255, 92, 58, 28))
    draw.ellipse((1180, -80, 1660, 400), fill=(255, 154, 115, 20))

    card = (70, 70, w - 70, h - 70)
    rounded(draw, card, 40, (12, 12, 12, 235), outline=(255, 255, 255, 26), width=2)

    logo = Image.open(LOGO).convert("RGBA")
    logo.thumbnail((210, 60))
    image.alpha_composite(logo, (110, 105))

    f_small = load_font(22, True)
    f_title = load_font(58, True)
    f_body = load_font(26, False)
    f_card_title = load_font(20, True)
    f_num = load_font(64, True)
    f_label = load_font(18, True)

    draw.text((110, 180), "STABLE - AI Reliability Engine", font=f_small, fill=ACCENT_SOFT)
    draw.text((110, 220), "Detect when your AI gives wrong or risky answers", font=f_title, fill=TEXT)
    draw.text((110, 315), "Generate an answer, run a second-pass evaluator, and surface the exact phrase that may be unreliable.", font=f_body, fill=MUTED)

    chip_y = 390
    chips = ["Response verification", "Hallucination review", "Risk classification"]
    x = 110
    for chip in chips:
        tw = draw.textbbox((0, 0), chip, font=f_small)[2]
        rounded(draw, (x, chip_y, x + tw + 38, chip_y + 46), 23, (255, 255, 255, 10), outline=(255,255,255,24), width=2)
        draw.text((x + 18, chip_y + 10), chip, font=f_small, fill=(248,221,209,255))
        x += tw + 58

    prompt = (110, 470, 1030, 790)
    rounded(draw, prompt, 28, (15, 15, 15, 255), outline=(255,255,255,20), width=2)
    draw.text((138, 500), "PROMPT", font=f_label, fill=ACCENT_SOFT)
    inner = (138, 545, 1002, 700)
    rounded(draw, inner, 24, (8, 8, 8, 255), outline=(255,255,255,18), width=2)
    draw.text((168, 585), "What is the capital of Australia?", font=load_font(32, False), fill=(255,245,240,255))

    buttons = [
        ("Analyze", True),
        ("Test hallucination", False),
        ("Test unsafe output", False),
    ]
    bx = 138
    by = 726
    for label, primary in buttons:
        tw = draw.textbbox((0,0), label, font=f_small)[2]
        bw = tw + 42
        fill = (255, 112, 74, 255) if primary else (255, 255, 255, 12)
        outline = None if primary else (255,255,255,26)
        rounded(draw, (bx, by, bx + bw, by + 52), 26, fill, outline=outline, width=2)
        txt_fill = (38, 10, 4, 255) if primary else (248,221,209,255)
        draw.text((bx + 20, by + 12), label, font=f_small, fill=txt_fill)
        bx += bw + 18

    side = (1070, 170, 1450, 790)
    rounded(draw, side, 30, (16, 12, 12, 255), outline=(255,255,255,22), width=2)
    draw.text((1105, 205), "LIVE ANALYSIS", font=f_label, fill=ACCENT_SOFT)

    sphere_center = (1260, 330)
    draw.ellipse((1165, 235, 1355, 425), fill=(255, 96, 61, 255))
    draw.ellipse((1140, 210, 1380, 450), outline=(255, 179, 146, 80), width=3)
    draw.ellipse((1115, 185, 1405, 475), outline=(255, 220, 200, 35), width=2)
    draw.text((1120, 470), "Confidence Score", font=f_card_title, fill=(255,174,144,255))
    draw.text((1118, 500), "81", font=f_num, fill=(255,244,239,255))

    risk_box = (1105, 610, 1415, 745)
    rounded(draw, risk_box, 22, (255,255,255,10), outline=(255,255,255,18), width=2)
    draw.text((1130, 635), "High Risk", font=load_font(30, True), fill=(255,210,196,255))
    draw.text((1130, 680), "Flagged phrase: Sydney", font=load_font(24, False), fill=(255,155,125,255))
    draw.text((1130, 713), "Correct capital is Canberra", font=load_font(21, False), fill=(209,178,164,255))

    image.convert("RGB").save(PRODUCT_IMG, quality=95)


def draw_background(c, width, height):
    c.setFillColor(HexColor(BG))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    c.setStrokeColor(Color(1, 1, 1, alpha=0.06))
    c.setLineWidth(0.4)
    step = 48
    for x in range(0, int(width), step):
        c.line(x, 0, x, height)
    for y in range(0, int(height), step):
        c.line(0, y, width, y)
    c.setFillColor(Color(1, 0.36, 0.22, alpha=0.10))
    c.circle(width * 0.12, height * 0.92, 110, fill=1, stroke=0)
    c.setFillColor(Color(1, 0.58, 0.45, alpha=0.08))
    c.circle(width * 0.88, height * 0.88, 90, fill=1, stroke=0)


def draw_logo(c, width, height):
    img = ImageReader(str(LOGO))
    iw, ih = img.getSize()
    target_w = 135
    target_h = target_w * ih / iw
    c.drawImage(img, 46, height - 56 - target_h, width=target_w, height=target_h, mask='auto')


def draw_tag(c, tag, x, y):
    c.setFillColor(Color(1, 0.36, 0.22, alpha=0.12))
    c.roundRect(x, y - 8, 190, 30, 14, fill=1, stroke=0)
    c.setFillColor(HexColor(ACCENT_SOFT))
    c.setFont("Helvetica-Bold", 10)
    c.drawString(x + 14, y + 2, tag)


def draw_stats(c, stats, x, y, width):
    box_w = (width - 20) / 3
    for i, (title, body) in enumerate(stats):
        bx = x + i * (box_w + 10)
        c.setFillColor(HexColor(CARD_ALT))
        c.roundRect(bx, y, box_w, 92, 16, fill=1, stroke=0)
        c.setFillColor(Color(1, 0.6, 0.46, alpha=0.16))
        c.roundRect(bx, y, box_w, 92, 16, fill=0, stroke=1)
        c.setFillColor(HexColor(ACCENT_SOFT))
        c.setFont("Helvetica-Bold", 12)
        c.drawString(bx + 14, y + 64, title)
        c.setFillColor(HexColor(TEXT))
        c.setFont("Helvetica", 11)
        ty = y + 45
        for line in wrap(body, 20):
            c.drawString(bx + 14, ty, line)
            ty -= 13


def draw_slide(c, slide, width, height):
    draw_background(c, width, height)
    draw_logo(c, width, height)
    draw_tag(c, slide["tag"], 46, height - 112)

    c.setFillColor(HexColor(TEXT))
    title_lines = wrap(slide["title"], 28)
    y = height - 170
    title_font = 28 if len(title_lines) > 2 else 32
    c.setFont("Helvetica-Bold", title_font)
    for line in title_lines:
        c.drawString(46, y, line)
        y -= title_font + 6

    c.setFillColor(HexColor(MUTED))
    c.setFont("Helvetica", 15)
    for line in wrap(slide["subtitle"], 65):
        c.drawString(46, y - 8, line)
        y -= 22

    if slide.get("image"):
        img_x, img_y, img_w, img_h = 46, 86, width - 92, y - 98
        c.setFillColor(HexColor(CARD))
        c.roundRect(img_x, img_y, img_w, img_h, 20, fill=1, stroke=0)
        c.drawImage(slide["image"], img_x + 10, img_y + 10, width=img_w - 20, height=img_h - 20, preserveAspectRatio=True, mask='auto')
        return

    card_x = 46
    card_y = 160 if slide.get("stats") else 74
    card_w = width - 92
    card_h = y - card_y - 12
    c.setFillColor(HexColor(CARD))
    c.roundRect(card_x, card_y, card_w, card_h, 20, fill=1, stroke=0)
    c.setFillColor(Color(1, 1, 1, alpha=0.06))
    c.roundRect(card_x, card_y, card_w, card_h, 20, fill=0, stroke=1)

    if slide.get("stats"):
        draw_stats(c, slide["stats"], 46, 74, width - 92)

    bullet_y = y - 44
    c.setFillColor(HexColor(TEXT))
    c.setFont("Helvetica-Bold", 16)
    c.drawString(card_x + 24, bullet_y, "Key Points")
    bullet_y -= 28
    c.setFont("Helvetica", 15)
    for bullet in slide["bullets"]:
        lines = wrap(bullet, 58)
        c.setFillColor(HexColor(ACCENT))
        c.circle(card_x + 28, bullet_y + 5, 3, fill=1, stroke=0)
        c.setFillColor(HexColor(TEXT))
        text_y = bullet_y
        for line in lines:
            c.drawString(card_x + 42, text_y, line)
            text_y -= 19
        bullet_y = text_y - 10


def rgb(hex_color):
    val = hex_color.replace("#", "")
    return RGBColor(int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    p.alignment = PP_ALIGN.LEFT
    return box


def add_round_rect(slide, x, y, w, h, fill, line="#FFFFFF", line_alpha=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    if line_alpha is not None:
        shape.line.transparency = line_alpha
    return shape


def add_stats_ppt(slide, stats):
    x0 = 0.5
    box_w = 4.08
    gap = 0.12
    for i, (title, body) in enumerate(stats):
        x = Inches(x0 + i * (box_w + gap))
        add_round_rect(slide, x, Inches(3.62), Inches(box_w), Inches(1.06), CARD_ALT, line=ACCENT_SOFT, line_alpha=0.55)
        add_textbox(slide, x + Inches(0.16), Inches(3.82), Inches(1.8), Inches(0.2), title, 12, True, ACCENT_SOFT)
        add_textbox(slide, x + Inches(0.16), Inches(4.12), Inches(box_w - 0.3), Inches(0.4), body, 11, False, TEXT)


def make_pdf():
    width, height = landscape((1280, 720))
    c = canvas.Canvas(str(PDF_PATH), pagesize=(width, height))
    for slide in slides:
        draw_slide(c, slide, width, height)
        c.showPage()
    c.save()


def make_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for slide_data in slides:
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(BG)

        for x in [0.2, 1.0, 1.8, 2.6, 3.4, 4.2, 5.0, 5.8, 6.6, 7.4, 8.2, 9.0, 9.8, 10.6, 11.4, 12.2]:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), 0, Pt(0.5), Inches(7.5))
            line.fill.solid(); line.fill.fore_color.rgb = rgb("#1A1A1A"); line.fill.transparency = 0.6
            line.line.fill.background()
        for y in [0.2, 0.9, 1.6, 2.3, 3.0, 3.7, 4.4, 5.1, 5.8, 6.5, 7.2]:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(y), Inches(13.333), Pt(0.5))
            line.fill.solid(); line.fill.fore_color.rgb = rgb("#1A1A1A"); line.fill.transparency = 0.6
            line.line.fill.background()

        slide.shapes.add_picture(str(LOGO), Inches(0.5), Inches(0.35), width=Inches(1.55))
        tag = add_round_rect(slide, Inches(0.5), Inches(1.12), Inches(2.25), Inches(0.38), "#2A1611", line=ACCENT_SOFT, line_alpha=0.65)
        tag.fill.transparency = 0.18
        add_textbox(slide, Inches(0.65), Inches(1.17), Inches(2.1), Inches(0.25), slide_data["tag"], 10, True, ACCENT_SOFT)

        add_textbox(slide, Inches(0.5), Inches(1.58), Inches(6.9), Inches(1.5), slide_data["title"], 28, True, TEXT)
        add_textbox(slide, Inches(0.5), Inches(2.82), Inches(7.0), Inches(0.95), slide_data["subtitle"], 14, False, MUTED)

        if slide_data.get("image"):
            add_round_rect(slide, Inches(0.5), Inches(3.72), Inches(12.33), Inches(2.95), CARD, line=BORDER, line_alpha=0.4)
            slide.shapes.add_picture(slide_data["image"], Inches(0.58), Inches(3.8), width=Inches(12.17), height=Inches(2.78))
            continue

        if slide_data.get("stats"):
            add_stats_ppt(slide, slide_data["stats"])
            card_y = Inches(4.86)
            card_h = Inches(1.82)
        else:
            card_y = Inches(3.78)
            card_h = Inches(2.92)

        add_round_rect(slide, Inches(0.5), card_y, Inches(12.33), card_h, CARD, line=BORDER, line_alpha=0.42)
        add_textbox(slide, Inches(0.78), card_y + Inches(0.22), Inches(2.0), Inches(0.2), "Key Points", 16, True, TEXT)

        y = float(card_y.inches) + 0.56
        for bullet in slide_data["bullets"]:
            lines = wrap(bullet, 54)
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.82), Inches(y + 0.06), Inches(0.08), Inches(0.08))
            dot.fill.solid(); dot.fill.fore_color.rgb = rgb(ACCENT); dot.line.fill.background()
            add_textbox(slide, Inches(1.0), Inches(y), Inches(10.9), Inches(0.7 + max(0, len(lines)-1)*0.18), "\n".join(lines), 14, False, TEXT)
            y += 0.46 + max(0, len(lines)-1) * 0.18

    prs.save(str(PPTX_PATH))


create_product_snapshot()
make_pdf()
make_pptx()
print(PDF_PATH)
print(PPTX_PATH)
print(PRODUCT_IMG)
